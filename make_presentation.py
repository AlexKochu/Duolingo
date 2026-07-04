import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6] # Blank layout

    # Colors
    c_dark_green = RGBColor(20, 60, 50)
    c_light_green = RGBColor(47, 93, 80)
    c_orange = RGBColor(255, 107, 74)
    c_dark_text = RGBColor(30, 35, 45)
    c_light_text = RGBColor(92, 107, 100)
    c_bg_light = RGBColor(245, 247, 246)
    c_white = RGBColor(255, 255, 255)
    c_border = RGBColor(220, 225, 222)
    c_danger_bg = RGBColor(253, 240, 238)
    
    # Helper to draw full background
    def set_bg(slide, color):
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        rect.fill.solid()
        rect.fill.fore_color.rgb = color
        rect.line.fill.background() # No border
        # Send to back by moving it to the first position in the shape list
        slide.shapes._spTree.remove(rect._element)
        slide.shapes._spTree.insert(2, rect._element)
        return rect

    # Helper to add standard text box
    def add_textbox(slide, left, top, width, height, text, size=14, bold=False, color=c_dark_text, align=PP_ALIGN.LEFT, italic=False):
        tx_box = slide.shapes.add_textbox(left, top, width, height)
        tf = tx_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = align
        p.font.name = 'Arial'
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.italic = italic
        p.font.color.rgb = color
        return tx_box

    # Helper to add formatted title/header on content slides
    def add_slide_header(slide, header_text, title_text):
        set_bg(slide, c_bg_light)
        # Header (gray small text)
        add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.3), header_text.upper(), size=11, bold=True, color=c_light_green)
        # Title (large text)
        add_textbox(slide, Inches(0.8), Inches(0.8), Inches(11.7), Inches(0.7), title_text, size=24, bold=True, color=c_dark_green)

    # Helper to draw a card shape (rectangle)
    def add_card(slide, left, top, width, height, bg_color=c_white, border_color=c_border):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        if border_color:
            card.line.color.rgb = border_color
            card.line.width = Pt(1.5)
        else:
            card.line.fill.background()
        return card

    # Helper to add bullet point list inside text box
    def add_bullets(slide, left, top, width, height, bullets, size=13, color=c_dark_text, spaced=True):
        tx_box = slide.shapes.add_textbox(left, top, width, height)
        tf = tx_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        for i, b in enumerate(bullets):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.text = "•  " + b
            p.font.name = 'Arial'
            p.font.size = Pt(size)
            p.font.color.rgb = color
            if spaced:
                p.space_after = Pt(8)
        return tx_box

    # Helper to add dynamic banner at the bottom
    def add_bottom_banner(slide, text, bg_color=c_dark_green, text_color=c_white, height=Inches(0.8), y_pos=Inches(6.2)):
        banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), y_pos, Inches(11.733), height)
        banner.fill.solid()
        banner.fill.fore_color.rgb = bg_color
        banner.line.fill.background()
        
        # Text inside
        tf = banner.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = PP_ALIGN.CENTER
        p.font.name = 'Arial'
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = text_color
        return banner

    # ==========================================
    # SLIDE 1: Title & Problem Statement (Dark Green Theme)
    # ==========================================
    s1 = prs.slides.add_slide(blank_layout)
    set_bg(s1, c_dark_green)
    
    # Subtitle Top Header
    add_textbox(s1, Inches(1.0), Inches(1.2), Inches(11.3), Inches(0.4), 
                "PRODUCT MANAGEMENT CASE STUDY  ·  VEDANTU PM INTERNSHIP", 
                size=12, bold=True, color=c_orange)
    
    # Main Title
    add_textbox(s1, Inches(1.0), Inches(1.7), Inches(11.3), Inches(1.0), 
                "Progressive Trust Verification", 
                size=44, bold=True, color=c_white)
    
    # Subtitle Description
    add_textbox(s1, Inches(1.0), Inches(2.7), Inches(11.3), Inches(0.5), 
                "Designing consent for Duolingo's under-18 users under India's DPDP Act, 2023", 
                size=18, color=RGBColor(190, 210, 200))
    
    # Problem Box (Light translucent box or clean light shape)
    p_box = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(3.6), Inches(7.5), Inches(1.8))
    p_box.fill.solid()
    p_box.fill.fore_color.rgb = RGBColor(28, 78, 66)
    p_box.line.color.rgb = c_light_green
    p_box.line.width = Pt(1.5)
    
    add_textbox(s1, Inches(1.3), Inches(3.8), Inches(6.9), Inches(0.4), "THE PROBLEM", size=11, bold=True, color=c_orange)
    add_textbox(s1, Inches(1.3), Inches(4.2), Inches(6.9), Inches(1.0), 
                "On a child's account, the person using the product and the person who must consent to its data practices are two different people. Duolingo's current flow is built for one user — not two.", 
                size=13.5, color=c_white)
    
    # Opportunity Box
    o_box = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(5.6), Inches(7.5), Inches(1.0))
    o_box.fill.solid()
    o_box.fill.fore_color.rgb = c_light_green
    o_box.line.fill.background()
    
    add_textbox(s1, Inches(1.3), Inches(5.75), Inches(6.9), Inches(0.3), "THE OPPORTUNITY", size=10, bold=True, color=c_orange)
    add_textbox(s1, Inches(1.3), Inches(6.05), Inches(6.9), Inches(0.5), 
                "Design a consent experience that protects children while preserving learning engagement.", 
                size=13, bold=True, color=c_white)

    # Aisha / Priya Visual Side Panel
    panel = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.2), Inches(1.7), Inches(3.1), Inches(4.9))
    panel.fill.solid()
    panel.fill.fore_color.rgb = c_white
    panel.line.fill.background()
    
    add_textbox(s1, Inches(9.4), Inches(2.2), Inches(2.7), Inches(0.6), "👤 Aisha, 12", size=18, bold=True, color=c_dark_green, align=PP_ALIGN.CENTER)
    add_textbox(s1, Inches(9.4), Inches(2.8), Inches(2.7), Inches(0.4), "Account Holder", size=11.5, color=c_light_text, align=PP_ALIGN.CENTER)
    
    add_textbox(s1, Inches(9.4), Inches(3.6), Inches(2.7), Inches(0.6), "≠", size=28, bold=True, color=c_orange, align=PP_ALIGN.CENTER)
    
    add_textbox(s1, Inches(9.4), Inches(4.4), Inches(2.7), Inches(0.6), "🛡️ Priya (Mother)", size=18, bold=True, color=c_dark_green, align=PP_ALIGN.CENTER)
    add_textbox(s1, Inches(9.4), Inches(5.0), Inches(2.7), Inches(0.4), "Consent Provider", size=11.5, color=c_light_text, align=PP_ALIGN.CENTER)
    
    add_textbox(s1, Inches(9.4), Inches(5.8), Inches(2.7), Inches(0.6), "One product. Two people to design for.", size=12, italic=True, color=c_light_green, align=PP_ALIGN.CENTER)
    
    add_textbox(s1, Inches(1.0), Inches(6.9), Inches(5.0), Inches(0.3), "Alex Kochumon Koshy  ·  July 2026", size=11, color=RGBColor(160, 180, 170))

    # ==========================================
    # SLIDE 2: Why Duolingo (Light Theme)
    # ==========================================
    s2 = prs.slides.add_slide(blank_layout)
    add_slide_header(s2, "Why This Product", "Most apps have one user. Duolingo has two.")
    add_textbox(s2, Inches(0.8), Inches(1.5), Inches(11.7), Inches(0.5), 
                "Consent design only gets hard when the person using the product isn't the person who can legally agree to its data practices.", 
                size=15, color=c_light_text)

    # Column 1: Spotify
    add_card(s2, Inches(0.8), Inches(2.2), Inches(5.6), Inches(3.7))
    add_textbox(s2, Inches(1.2), Inches(2.5), Inches(4.8), Inches(0.4), "🎵 Spotify (Baseline Case)", size=18, bold=True, color=c_dark_green)
    add_textbox(s2, Inches(1.2), Inches(3.0), Inches(4.8), Inches(0.3), "User = Consent Provider", size=13, bold=True, color=c_orange)
    spotify_bullets = [
        "The adult who signs up is the same person whose data is collected.",
        "One consent flow, one screen, one moment — it just works.",
        "Zero friction handoff since identity is single-layered."
    ]
    add_bullets(s2, Inches(1.2), Inches(3.5), Inches(4.8), Inches(2.0), spotify_bullets)

    # Column 2: Duolingo
    add_card(s2, Inches(6.9), Inches(2.2), Inches(5.6), Inches(3.7))
    add_textbox(s2, Inches(7.3), Inches(2.5), Inches(4.8), Inches(0.4), "🦉 Duolingo Child Accounts (Divergent)", size=18, bold=True, color=c_dark_green)
    add_textbox(s2, Inches(7.3), Inches(3.0), Inches(4.8), Inches(0.3), "User ≠ Consent Provider", size=13, bold=True, color=c_orange)
    duolingo_bullets = [
        "Aisha (12) uses the product daily — she is the data subject.",
        "Priya (her mother) must consent on her behalf — she never opens the app.",
        "Consent request is detached from the moment of actual usage."
    ]
    add_bullets(s2, Inches(7.3), Inches(3.5), Inches(4.8), Inches(2.0), duolingo_bullets)

    add_bottom_banner(s2, "So what: A design built for the Spotify case will always under-serve the parent — and eventually lose the child's account too.")

    # ==========================================
    # SLIDE 3: Personas (Light Theme)
    # ==========================================
    s3 = prs.slides.add_slide(blank_layout)
    add_slide_header(s3, "Who We're Designing For", "Two personas, one account, pulling in different directions")

    # Column 1: Aisha
    add_card(s3, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.1))
    add_textbox(s3, Inches(1.2), Inches(2.1), Inches(4.8), Inches(0.4), "👩‍🎓 Aisha, 12 (Account Holder)", size=18, bold=True, color=c_dark_green)
    add_textbox(s3, Inches(1.2), Inches(2.6), Inches(4.8), Inches(0.3), "Wants speed, momentum, and instant learning feedback", size=12, italic=True, color=c_light_text)
    
    aisha_bullets = [
        "Wants to learn Spanish for school and keep her daily streak alive.",
        "She won't read a consent screen — she will tap through anything standing between her and the next learning unit.",
        "Highly sensitive to UX friction and layout interruptions."
    ]
    add_bullets(s3, Inches(1.2), Inches(3.1), Inches(4.8), Inches(2.4), aisha_bullets)

    # Column 2: Priya
    add_card(s3, Inches(6.9), Inches(1.8), Inches(5.6), Inches(4.1))
    add_textbox(s3, Inches(7.3), Inches(2.1), Inches(4.8), Inches(0.4), "🛡️ Priya (Aisha's Mother — Consent Provider)", size=18, bold=True, color=c_dark_green)
    add_textbox(s3, Inches(7.3), Inches(2.6), Inches(4.8), Inches(0.3), "Wants absolute safety, visibility, and control over data", size=12, italic=True, color=c_light_text)
    
    priya_bullets = [
        "Has zero default visibility today into what data is being collected.",
        "Cares about: Is her child's voice being recorded? Is it shared with third-party advertisers? How does she revoke permission?",
        "Refuses to read lengthy Privacy Policy documents to find simple answers."
    ]
    add_bullets(s3, Inches(7.3), Inches(3.1), Inches(4.8), Inches(2.4), priya_bullets)

    add_bottom_banner(s3, "The tension: The product's UX is optimized for a 12-year-old's engagement. What her mother needs is visibility and control. Those goals pull in opposite directions.")

    # ==========================================
    # SLIDE 4: Current Experience (Light Theme)
    # ==========================================
    s4 = prs.slides.add_slide(blank_layout)
    add_slide_header(s4, "Current Experience", "Today's flow creates data risk with no visibility for parents")

    # 2x2 Grid of cards
    grid_data = [
        ("🎂 Self-declared birthdate", "Captured at signup with no verification step. Creates non-compliance risk if minor accounts bypass locks.", Inches(0.8), Inches(1.8)),
        ("🎙️ Voice-based exercises", "Speaking lessons require microphone access. The device prompts generic OS popup with no parental alignment.", Inches(6.9), Inches(1.8)),
        ("📈 Social & Leaderboards", "User progress is shared on public boards. Behavioral metrics inform recommendations without parent setup.", Inches(0.8), Inches(3.9)),
        ("🚫 No parent-facing view", "Parents have no dashboard or panel to audit settings, view stored details, or easily withdraw consent.", Inches(6.9), Inches(3.9))
    ]

    for title, desc, x, y in grid_data:
        add_card(s4, x, y, Inches(5.6), Inches(1.8))
        add_textbox(s4, x + Inches(0.4), y + Inches(0.2), Inches(4.8), Inches(0.4), title, size=16, bold=True, color=c_dark_green)
        add_textbox(s4, x + Inches(0.4), y + Inches(0.7), Inches(4.8), Inches(0.9), desc, size=13, color=c_dark_text)

    add_bottom_banner(s4, "The trust gap: A parent who discovers this after the fact has no reason to trust the product going forward — that's a retention risk, independent of any regulation.")

    # ==========================================
    # SLIDE 5: Prototype Walkthrough (Light Theme)
    # ==========================================
    s5 = prs.slides.add_slide(blank_layout)
    add_slide_header(s5, "Prototype Walkthrough", "Real-time progressive feature activation flow")
    
    # Process sequence flow text
    add_textbox(s5, Inches(0.8), Inches(1.5), Inches(11.7), Inches(0.5),
                "Safe Mode  ⟶  Risk Trigger  ⟶  Parent Review  ⟶  Feature Unlock",
                size=16, bold=True, color=c_orange, align=PP_ALIGN.CENTER)

    # 3 screenshot placeholder cards side-by-side
    s_col_w = Inches(3.6)
    s_col_gap = Inches(0.4)
    s_left_margin = Inches(0.8)
    
    screenshots = [
        ("📱 SCREENSHOT 1: Safe Mode", 
         "Child can immediately start learning while sensitive features remain locked.", 
         "Aisha can do reading/listening lessons instantly without roadblocks. Mic & leaderboard elements display lock (🔒) icons."),
        ("📩 SCREENSHOT 2: Parent Review", 
         "Parent receives contextual explanation and granular controls.", 
         "Priya opens family link showing why Aisha needs the mic (speaking practice) with toggle controls to customize data sharing."),
        ("🎉 SCREENSHOT 3: Approved Screen", 
         "Child receives immediate value once permission is approved.", 
         "App syncs live with a confetti burst. Microphone unlocks and Aisha gains +15 XP from speaking exercises.")
    ]

    for i, (title, desc, detail) in enumerate(screenshots):
        x = s_left_margin + i * (s_col_w + s_col_gap)
        add_card(s5, x, Inches(2.2), s_col_w, Inches(3.7))
        # Draw placeholder for screenshot image
        rect_img = s5.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(0.3), Inches(2.4), Inches(3.0), Inches(1.8))
        rect_img.fill.solid()
        rect_img.fill.fore_color.rgb = c_bg_light
        rect_img.line.color.rgb = c_border
        
        # Label inside placeholder
        tf = rect_img.text_frame
        p = tf.paragraphs[0]
        p.text = "[ Screenshot Placeholder ]"
        p.alignment = PP_ALIGN.CENTER
        p.font.name = 'Arial'
        p.font.size = Pt(11)
        p.font.italic = True
        p.font.color.rgb = c_light_text

        # Text below
        add_textbox(s5, x + Inches(0.2), Inches(4.3), Inches(3.2), Inches(0.4), title, size=13, bold=True, color=c_orange)
        add_textbox(s5, x + Inches(0.2), Inches(4.7), Inches(3.2), Inches(0.5), desc, size=11.5, bold=True, color=c_dark_green)
        add_textbox(s5, x + Inches(0.2), Inches(5.2), Inches(3.2), Inches(0.6), detail, size=11, color=c_dark_text)

    # ==========================================
    # SLIDE 6: PTV Solution (Light Theme)
    # ==========================================
    s6 = prs.slides.add_slide(blank_layout)
    add_slide_header(s6, "The Solution", "Progressive Trust Verification (PTV)")

    # Accent Header Callout
    acc = s6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.6), Inches(11.733), Inches(0.6))
    acc.fill.solid()
    acc.fill.fore_color.rgb = c_danger_bg
    acc.line.color.rgb = c_orange
    
    add_textbox(s6, Inches(1.0), Inches(1.75), Inches(11.3), Inches(0.4),
                "Core Principle: Ask for parental consent when data risk appears — not when the account is created.",
                size=14, bold=True, color=c_orange)

    # Horizontal step blocks
    step_w = Inches(2.1)
    step_gap = Inches(0.3)
    steps = [
        ("1. Signup", "Birthdate self-declared. No upfront parental block."),
        ("2. Safe Mode", "Default state. No voice, ads, or profile sharing."),
        ("3. Risk Trigger", "First time voice/location is requested by lesson."),
        ("4. Parent Review", "Priya reviews granular controls via cloud link."),
        ("5. Sync Live", "Live update unlocks feature with sound & XP.")
    ]

    for i, (title, desc) in enumerate(steps):
        x = Inches(0.8) + i * (step_w + step_gap)
        add_card(s6, x, Inches(2.4), step_w, Inches(2.2))
        add_textbox(s6, x + Inches(0.15), Inches(2.6), step_w - Inches(0.3), Inches(0.4), title, size=14, bold=True, color=c_dark_green)
        add_textbox(s6, x + Inches(0.15), Inches(3.1), step_w - Inches(0.3), Inches(1.3), desc, size=12, color=c_dark_text)

    add_bottom_banner(s6, "Why this beats a signup wall: A wall at signup protects data the account doesn't have yet — at the cost of activation. PTV lets the child get real value first, and only asks the parent to engage at the exact moment there's something concrete worth reviewing.")

    # ==========================================
    # SLIDE 7: Metrics + Success Criteria (Light Theme)
    # ==========================================
    s7 = prs.slides.add_slide(blank_layout)
    add_slide_header(s7, "How We'll Know It's Working", "Metrics that tie trust to business health")

    # Column 1: North star card (Large card)
    add_card(s7, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.9), bg_color=c_dark_green, border_color=None)
    add_textbox(s7, Inches(1.2), Inches(2.1), Inches(4.8), Inches(0.4), "🌟 NORTH STAR METRIC", size=13, bold=True, color=c_orange)
    
    metric_val = "% of under-18 accounts with valid, current, feature-level data-collecting feature in active use."
    add_textbox(s7, Inches(1.2), Inches(2.6), Inches(4.8), Inches(1.8), metric_val, size=24, bold=True, color=c_white)
    
    metric_desc = "This is the single number that tells us the product is both legally compliant (DPDP Act verified) and parent-trusted (permissions active) at the same time without dampening child user activity."
    add_textbox(s7, Inches(1.2), Inches(4.6), Inches(4.8), Inches(1.8), metric_desc, size=13.5, color=c_white)

    # Column 2: Supporting Metrics & Guardrails
    add_card(s7, Inches(6.9), Inches(1.8), Inches(5.6), Inches(4.9))
    
    add_textbox(s7, Inches(7.3), Inches(2.1), Inches(4.8), Inches(0.4), "📋 SUPPORTING PRODUCT METRICS", size=14, bold=True, color=c_dark_green)
    supporting_metrics = [
        "Time-to-Verification: Trigger fired ⟶ parent completes cloud review.",
        "Safe Mode Adoption: Engagement levels of children without voice/social active.",
        "Consent Revocation Rate: Percentage of parents pulling back permissions."
    ]
    add_bullets(s7, Inches(7.3), Inches(2.6), Inches(4.8), Inches(1.8), supporting_metrics, size=12.5)

    # Line separator
    line = s7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.3), Inches(4.4), Inches(4.8), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = c_border
    line.line.fill.background()

    add_textbox(s7, Inches(7.3), Inches(4.6), Inches(4.8), Inches(0.4), "🚨 BUSINESS GUARDRAILS", size=14, bold=True, color=c_dark_green)
    guardrail_metrics = [
        "D7 / D30 Cohort Retention: Ensures trust mechanics do not negatively impact learning streak frequency.",
        "Consent Support Tickets: Tracking parental complaints as a UX friction benchmark."
    ]
    add_bullets(s7, Inches(7.3), Inches(5.1), Inches(4.8), Inches(1.5), guardrail_metrics, size=12.5)

    # ==========================================
    # SLIDE 8: Diagnostic Thinking & Dashboard (Light Theme)
    # ==========================================
    s8 = prs.slides.add_slide(blank_layout)
    add_slide_header(s8, "Diagnostic Thinking & Management", "PM Dashboard & Response Framework")

    # Column 1: Dashboard Design
    add_card(s8, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.9))
    add_textbox(s8, Inches(1.2), Inches(2.0), Inches(4.8), Inches(0.4), "📊 PM Dashboard Design", size=16, bold=True, color=c_dark_green)
    
    # KPI Grid (4 visual small boxes)
    kpis = [
        ("Consent Rate", "62%", Inches(1.2), Inches(2.5)),
        ("Pending", "18%", Inches(3.7), Inches(2.5)),
        ("Safe Mode", "15%", Inches(1.2), Inches(3.7)),
        ("Revoked", "5%", Inches(3.7), Inches(3.7))
    ]
    for label, val, x, y in kpis:
        box = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.2), Inches(1.0))
        box.fill.solid()
        box.fill.fore_color.rgb = c_bg_light
        box.line.color.rgb = c_border
        box.line.width = Pt(1)
        add_textbox(s8, x + Inches(0.15), y + Inches(0.15), Inches(1.9), Inches(0.3), label, size=11, color=c_light_text)
        add_textbox(s8, x + Inches(0.15), y + Inches(0.45), Inches(1.9), Inches(0.5), val, size=22, bold=True, color=c_orange)

    # Chart descriptions bottom
    add_textbox(s8, Inches(1.2), Inches(4.9), Inches(4.8), Inches(1.5),
                "Main Visual Charts:\n1. Consent Trend (Conversion rate of requests over time)\n2. Retention Trend (D30 streak correlation by consent state)\n3. Revocation Trend (Voluntary opt-out spikes matching digests)",
                size=12, color=c_dark_text)

    # Column 2: Diagnostic Thinking
    add_card(s8, Inches(6.9), Inches(1.8), Inches(5.6), Inches(4.9))
    add_textbox(s8, Inches(7.3), Inches(2.0), Inches(4.8), Inches(0.4), "🧠 Diagnostic Thinking (Triage Matrix)", size=16, bold=True, color=c_dark_green)
    add_textbox(s8, Inches(7.3), Inches(2.4), Inches(4.8), Inches(0.4), "If consent rate drops, systematically check:", size=12, color=c_light_text)

    diagnostics = [
        ("1. Notification Delivery Failure", "Check SMS/WhatsApp/Email webhook delivery rates. A drop here signals latency in parent communication paths."),
        ("2. Revocation Spike Correlation", "Analyze if parental opt-out spiked following the deployment of a quarterly email digest or billing cycle."),
        ("3. Cohort Dilution & Age Shift", "Validate if new user signups are skewing older/younger, changing the proportion of safe mode defaults.")
    ]

    dy = Inches(2.8)
    for title, desc in diagnostics:
        add_textbox(s8, Inches(7.3), dy, Inches(4.8), Inches(0.3), title, size=13, bold=True, color=c_orange)
        add_textbox(s8, Inches(7.3), dy + Inches(0.3), Inches(4.8), Inches(0.8), desc, size=11.5, color=c_dark_text)
        dy += Inches(1.1)

    # ==========================================
    # SLIDE 9: Rollout Plan (Light Theme)
    # ==========================================
    s9 = prs.slides.add_slide(blank_layout)
    add_slide_header(s9, "Getting This Live", "Rollout strategy and the risks we're accepting")

    # Three horizontal cards representing phases
    card_w = Inches(3.6)
    card_h = Inches(2.2)
    card_gap = Inches(0.4)
    left_margin = Inches(0.8)
    
    phases = [
        ("Phase 1 (1-2 weeks)", "Shadow Mode", [
            "Log trigger actions on production databases.",
            "Confirm logic works in background without changing active child UI.",
            "Verify cloud messaging load testing."
        ]),
        ("Phase 2 (2-4 weeks)", "New Signups Only", [
            "Limit risk to new acquisitions.",
            "Isolate experimental group to evaluate onboarding signup conversion.",
            "Check for immediate churn effects."
        ]),
        ("Phase 3 (6+ weeks)", "Staggered Existing", [
            "Staged region rollout (India first for DPDP compliance check).",
            "Slow rollout to the retained active cohort.",
            "Monitor customer support tick volumes."
        ])
    ]

    for i, (p_title, p_name, p_bullets) in enumerate(phases):
        x = left_margin + i * (card_w + card_gap)
        add_card(s9, x, Inches(1.8), card_w, card_h)
        add_textbox(s9, x + Inches(0.3), Inches(2.0), card_w - Inches(0.6), Inches(0.3), p_title, size=12, bold=True, color=c_orange)
        add_textbox(s9, x + Inches(0.3), Inches(2.3), card_w - Inches(0.6), Inches(0.4), p_name, size=16, bold=True, color=c_dark_green)
        add_bullets(s9, x + Inches(0.3), Inches(2.8), card_w - Inches(0.6), Inches(1.1), p_bullets, size=11, spaced=False)

    # Bottom guardrails
    g_box = s9.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(4.2), Inches(11.733), Inches(0.7))
    g_box.fill.solid()
    g_box.fill.fore_color.rgb = c_bg_light
    g_box.line.color.rgb = c_border
    
    add_textbox(s9, Inches(1.0), Inches(4.3), Inches(11.3), Inches(0.5),
                "Rollout Guardrails: Validate metrics daily. Immediately pause rollout if D30 cohort retention drops by >5%. A/B test UX details (SMS copy, button text) — never compliance necessity.",
                size=12, italic=True, color=c_light_green)

    # Risks list below
    add_textbox(s9, Inches(0.8), Inches(5.1), Inches(11.7), Inches(0.3), "KEY RISKS TO MANAGE", size=13, bold=True, color=c_dark_green)
    
    risks_y = Inches(5.4)
    r_cols = [
        ("🛠️ Engineering Complexity", "Every downstream microservice touching voice, social or ad cookies must query centralized consent state.", Inches(0.8)),
        ("📉 Monetization Slowdown", "Safe Mode limits target advertising & premium trials. Accept lower LTV during A/B evaluation.", Inches(4.8)),
        ("⏱️ Edge-Case Latency", "If parent is offline, child gets stuck in restricted features (mitigate via proactive Free Play suggestions).", Inches(8.8))
    ]
    for title, text, rx in r_cols:
        add_textbox(s9, rx, risks_y, Inches(3.7), Inches(0.3), title, size=13, bold=True, color=c_orange)
        add_textbox(s9, rx, risks_y + Inches(0.3), Inches(3.7), Inches(0.8), text, size=11, color=c_dark_text)

    # ==========================================
    # SLIDE 10: Risks & Tradeoffs (Light Theme)
    # ==========================================
    s10 = prs.slides.add_slide(blank_layout)
    add_slide_header(s10, "Product Tradeoffs & Key Takeaways", "Consent Wall vs. Progressive Trust Verification")

    # Add PowerPoint Table
    table_shape = s10.shapes.add_table(5, 3, Inches(0.8), Inches(1.8), Inches(11.733), Inches(3.6))
    table = table_shape.table
    
    # Column widths
    table.columns[0].width = Inches(1.8)
    table.columns[1].width = Inches(4.966)
    table.columns[2].width = Inches(4.966)

    # Table content matrix
    headers = ["Dimension", "Consent Wall (Signup Gate)", "Progressive Trust Verification (PTV)"]
    rows_data = [
        ["Compliance", "High — consent exists before any customer data begins to flow.", "High — consent is tied exactly to the specific feature in use."],
        ["Activation", "Low — heavy signup block causes high user drop-off or falsified age entry.", "Higher — child experiences product value first; parent active when needed."],
        ["Trust", "Fragile — one-time checklist checkbox; parent receives no ongoing relationship signals.", "Compounding — granular feature controls plus a recurring data digest report."],
        ["User Experience", "Interrupts the moment the child is most excited to set up and start learning.", "Interrupts only when the child chooses to unlock advanced features."]
    ]

    # Format Header Row
    for col_idx, header_text in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header_text
        cell.fill.solid()
        cell.fill.fore_color.rgb = c_dark_green
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.LEFT if col_idx == 0 else PP_ALIGN.CENTER
            p.font.name = 'Arial'
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = c_white

    # Format Data Rows
    for row_idx, row in enumerate(rows_data):
        for col_idx, text in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = text
            cell.fill.solid()
            # Alternate row backgrounds
            cell.fill.fore_color.rgb = c_white if row_idx % 2 == 0 else c_bg_light
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT
                p.font.name = 'Arial'
                p.font.size = Pt(12)
                p.font.color.rgb = c_dark_text
                if col_idx == 0:
                    p.font.bold = True
                    p.font.color.rgb = c_dark_green
                    
    # Bottom key takeaway banner
    add_bottom_banner(s10, "Key Takeaway: Progressive Trust Verification enables feature-level parental consent without introducing a high-friction onboarding experience.")

    # Save presentation
    output_path = "/Users/alexkochu/Desktop/Duolingo/Vedantu_Case_Study.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully to {output_path}")

if __name__ == '__main__':
    create_presentation()
